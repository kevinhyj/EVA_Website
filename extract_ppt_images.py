"""
Extract all images from slide 9 of RNAVerse_web.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH = r"D:\workspace\lingang_evaweb\RNAVerse_web.pptx"
OUTPUT_DIR = r"D:\workspace\lingang_evaweb\RNAVerse_Web\RNA-Verse\public\ppt"
SLIDE_INDEX = 4  # 0-based, slide 5

os.makedirs(OUTPUT_DIR, exist_ok=True)

prs = Presentation(PPTX_PATH)
slide = prs.slides[SLIDE_INDEX]

print(f"Slide 9 has {len(slide.shapes)} shapes")

count = 0

def extract_images(shapes, indent=0):
    global count
    prefix = "  " * indent
    for i, shape in enumerate(shapes):
        print(f"{prefix}Shape {i}: type={shape.shape_type}, name={shape.name}")
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            ext = image.ext
            fname = f"slide5_{count+1}.{ext}"
            fpath = os.path.join(OUTPUT_DIR, fname)
            with open(fpath, "wb") as f:
                f.write(image.blob)
            print(f"{prefix}  -> Saved: {fname} ({len(image.blob)} bytes)")
            count += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            extract_images(shape.shapes, indent + 1)

extract_images(slide.shapes)
print(f"\nDone. Saved {count} images to {OUTPUT_DIR}")
